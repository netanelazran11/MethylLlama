"""
Generate PowerPoint presentation: BMFM-DNA Initialized MethylLlama
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT      = RGBColor(0x16, 0x21, 0x3E)   # dark blue panel
HIGHLIGHT   = RGBColor(0x00, 0xB4, 0xD8)   # cyan
HIGHLIGHT2  = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
GREEN       = RGBColor(0x06, 0xD6, 0xA0)   # green for OK
ORANGE      = RGBColor(0xFF, 0xB7, 0x03)   # orange for numbers
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
SUBTITLE_GRAY = RGBColor(0xAA, 0xBB, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=17, title=None, title_size=19,
                   title_color=HIGHLIGHT, bullet_color=WHITE,
                   bullet="  •  ", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = bullet + item[0]
            run.font.color.rgb = item[1]
        else:
            run.text = bullet + item
            run.font.color.rgb = bullet_color
        run.font.size = Pt(font_size)
    return txBox


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)

# Top accent bar
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)

# Title
add_textbox(slide,
    "DNA-Informed Methylation Pretraining",
    Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
    font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide,
    "Initializing MethylLlama CpG Embeddings with BMFM-DNA Genomic Context",
    Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.8),
    font_size=22, color=HIGHLIGHT2, align=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4), Inches(3.75), Inches(5.3), Inches(0.04), HIGHLIGHT)

# Author / context
add_textbox(slide,
    "Netanel Azran  |  BMFM-RNA Methylation  |  Hebrew University",
    Inches(1), Inches(4.1), Inches(11.3), Inches(0.5),
    font_size=16, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)


# =============================================================================
# SLIDE 2 — The Problem
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "The Problem", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(4.5), Inches(0.04), HIGHLIGHT)

# Left panel — current approach
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.5), Inches(5.5), ACCENT)
add_textbox(slide, "Current Approach", Inches(0.6), Inches(1.3), Inches(5), Inches(0.5),
            font_size=20, bold=True, color=ORANGE)
add_textbox(slide,
    "Each CpG site (e.g. cg00000292) gets a random vector in the embedding table.\n\n"
    "The model has no idea that two CpGs living in the same promoter region are biologically similar — "
    "it must learn all CpG relationships from scratch, using only methylation values.\n\n"
    "This is slow and data-hungry.",
    Inches(0.6), Inches(1.9), Inches(5.1), Inches(4.2),
    font_size=16, color=LIGHT_GRAY)

# Right panel — insight
add_rect(slide, Inches(6.2), Inches(1.2), Inches(6.5), Inches(5.5), ACCENT)
add_textbox(slide, "Key Insight", Inches(6.4), Inches(1.3), Inches(6), Inches(0.5),
            font_size=20, bold=True, color=GREEN)
add_textbox(slide,
    "Every CpG site has a fixed genomic address — a specific location in the human genome.\n\n"
    "The DNA sequence surrounding that location encodes rich biological information:\n\n"
    "   • Is it in a CpG island?\n"
    "   • Is it in a gene promoter?\n"
    "   • Are there transcription factor binding motifs nearby?\n"
    "   • Is it in a repeat element?\n\n"
    "This information exists before seeing any methylation data — we can use it to give "
    "each CpG a biologically meaningful starting embedding.",
    Inches(6.4), Inches(1.9), Inches(6.1), Inches(4.2),
    font_size=16, color=LIGHT_GRAY)

# Arrow in middle
add_textbox(slide, "→", Inches(5.55), Inches(3.6), Inches(0.6), Inches(0.6),
            font_size=36, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 3 — What is a CpG Site
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Background: CpG Sites & DNA Methylation",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(7), Inches(0.04), HIGHLIGHT)

# Three boxes
for i, (title, body, col) in enumerate([
    ("What is a CpG Site?",
     "A position in the genome where Cytosine (C) is followed by Guanine (G). "
     "DNA methylation — adding a methyl group to the C — happens almost exclusively at CpG sites.",
     HIGHLIGHT),
    ("Beta Value",
     "The methylation level at a CpG, measured between 0 and 1.\n\n"
     "0 = completely unmethylated\n"
     "1 = completely methylated\n\n"
     "This is the core signal in the methylation array data.",
     ORANGE),
    ("Why 21,368 CpGs?",
     "The Illumina HM450 array measures ~450,000 CpG sites. "
     "AltumAge selected 21,368 'type-3' CpGs that are most informative for biological age prediction. "
     "These are the CpGs our model focuses on.",
     GREEN),
]):
    left = Inches(0.4 + i * 4.3)
    add_rect(slide, left, Inches(1.2), Inches(4.1), Inches(5.5), ACCENT)
    add_textbox(slide, title, left + Inches(0.2), Inches(1.35), Inches(3.7), Inches(0.5),
                font_size=18, bold=True, color=col)
    add_rect(slide, left + Inches(0.2), Inches(1.85), Inches(3.5), Inches(0.03), col)
    add_textbox(slide, body, left + Inches(0.2), Inches(1.95), Inches(3.7), Inches(4.4),
                font_size=15, color=LIGHT_GRAY)


# =============================================================================
# SLIDE 4 — The Data Downloads
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Required Data: What We Download and Why",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(7), Inches(0.04), HIGHLIGHT)

# Data item 1
add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(1.6), ACCENT)
add_textbox(slide, "1   Illumina HM450 hg38 Manifest", Inches(0.7), Inches(1.28),
            Inches(8), Inches(0.45), font_size=18, bold=True, color=ORANGE)
add_textbox(slide,
    "A table mapping every probe ID (e.g. cg00000292) to its exact chromosomal position "
    "(chromosome + base-pair coordinate) in the hg38 human genome. "
    "Without this, we cannot find where each CpG lives in the genome.",
    Inches(0.7), Inches(1.72), Inches(12), Inches(0.9), font_size=15, color=LIGHT_GRAY)

# Data item 2
add_rect(slide, Inches(0.4), Inches(3.0), Inches(12.5), Inches(1.8), ACCENT)
add_textbox(slide, "2   hg38 Reference Genome  (~3.2 GB)", Inches(0.7), Inches(3.08),
            Inches(8), Inches(0.45), font_size=18, bold=True, color=HIGHLIGHT2)
add_textbox(slide,
    "The complete human DNA sequence — all ~3 billion base pairs across every chromosome. "
    "Once we know a CpG's position from the manifest, we extract ±512 base pairs of DNA "
    "sequence centered on that CpG. This local window is what BMFM-DNA processes.",
    Inches(0.7), Inches(3.52), Inches(12), Inches(1.1), font_size=15, color=LIGHT_GRAY)

# Data item 3
add_rect(slide, Inches(0.4), Inches(4.95), Inches(12.5), Inches(1.8), ACCENT)
add_textbox(slide, "3   BMFM-DNA Model Checkpoint  (~1.4 GB)", Inches(0.7), Inches(5.03),
            Inches(8), Inches(0.45), font_size=18, bold=True, color=GREEN)
add_textbox(slide,
    "IBM Research's SCModernBert 113M model, pre-trained on GRCh38 (same as hg38). "
    "This model has learned to produce rich 768-dimensional representations of any DNA sequence. "
    "We use it as a frozen encoder — no fine-tuning needed.",
    Inches(0.7), Inches(5.47), Inches(12), Inches(1.1), font_size=15, color=LIGHT_GRAY)


# =============================================================================
# SLIDE 5 — Why BMFM-DNA not RNA
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Why BMFM-DNA and Not an RNA Model?",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(6), Inches(0.04), HIGHLIGHT)

# Table header
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.1), Inches(0.5), HIGHLIGHT)
add_rect(slide, Inches(6.6), Inches(1.2), Inches(6.3), Inches(0.5), RGBColor(0x80, 0x40, 0x40))
add_textbox(slide, "BMFM-DNA", Inches(0.4), Inches(1.22), Inches(6), Inches(0.45),
            font_size=19, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
add_textbox(slide, "RNA Model", Inches(6.6), Inches(1.22), Inches(6.2), Inches(0.45),
            font_size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Input", "DNA sequence  (A, C, G, T)", "RNA sequence  (A, C, G, U)"),
    ("Trained on", "GRCh38 reference genome", "Transcriptomes / RNA-seq"),
    ("Captures", "Regulatory elements, TF binding sites,\nCpG islands, repeat elements", "Splicing patterns, codon usage, UTRs"),
    ("Relevance to\nmethylation", "DIRECT — methylation happens ON the\nDNA, at the exact sequence level", "INDIRECT — RNA is downstream of\nmethylation, not the source"),
    ("Conclusion", "✓  Correct choice", "✗  Wrong level of biology"),
]

for i, (label, dna_text, rna_text) in enumerate(rows):
    y = Inches(1.75 + i * 1.0)
    bg = RGBColor(0x12, 0x18, 0x30) if i % 2 == 0 else ACCENT
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.95), bg)
    add_textbox(slide, label, Inches(0.5), y + Inches(0.05), Inches(1.5), Inches(0.85),
                font_size=13, bold=True, color=SUBTITLE_GRAY)
    col_dna = GREEN if i == 4 else LIGHT_GRAY
    col_rna = RGBColor(0xFF, 0x66, 0x66) if i == 4 else LIGHT_GRAY
    add_textbox(slide, dna_text, Inches(2.1), y + Inches(0.05), Inches(4.7), Inches(0.85),
                font_size=14, color=col_dna)
    add_textbox(slide, rna_text, Inches(7.0), y + Inches(0.05), Inches(5.7), Inches(0.85),
                font_size=14, color=col_rna)


# =============================================================================
# SLIDE 6 — Creating CpG Embeddings
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Step 1: Creating CpG Embeddings with BMFM-DNA",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=30, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.04), HIGHLIGHT)

add_textbox(slide, "One-time preprocessing  •  runs once  •  saves 21,368 × 768 vectors",
            Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
            font_size=14, color=SUBTITLE_GRAY, italic=True)

# Pipeline steps as boxes with arrows
steps = [
    ("Probe ID\ncg00000292", ORANGE),
    ("Manifest\nLookup", HIGHLIGHT),
    ("chr1 : 1,234,567", HIGHLIGHT2),
    ("Extract ±512 bp\nfrom hg38", GREEN),
    ("DNA Sequence\n1024 characters", HIGHLIGHT2),
    ("BMFM-DNA\n113M params", ORANGE),
    ("768-dim\nVector", GREEN),
]

box_w = Inches(1.55)
box_h = Inches(1.1)
gap   = Inches(0.25)
start_x = Inches(0.3)
y = Inches(1.7)

for i, (label, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(slide, x, y, box_w, box_h, ACCENT)
    # colored top bar
    add_rect(slide, x, y, box_w, Inches(0.06), color)
    add_textbox(slide, label, x + Inches(0.05), y + Inches(0.1),
                box_w - Inches(0.1), box_h - Inches(0.1),
                font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.03)
        add_textbox(slide, "→", ax, y + Inches(0.25), gap + Inches(0.1), Inches(0.6),
                    font_size=20, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Done for all 21368
add_textbox(slide,
    "Repeated for all 21,368 CpG sites  →  saved as  cpg_embeddings_bmfdna_21k.npy  [21368 × 768]",
    Inches(0.4), Inches(3.05), Inches(12.5), Inches(0.4),
    font_size=15, color=GREEN, bold=True)

# Three insight boxes
for i, (title, body) in enumerate([
    ("What the vector encodes",
     "The 768-dim vector summarizes everything about the local DNA context: "
     "CpG island membership, nearby regulatory motifs, chromatin accessibility signals learned from sequence."),
    ("Why pooler_output not mean pooling",
     "BMFM-DNA was trained with a learned pooling layer on the CLS token. "
     "Mean pooling collapses all CpGs to nearly identical vectors. "
     "Pooler output preserves biological structure (PC1 = 33.5% variance)."),
    ("Result",
     "Norm of each vector ≈ 7.62 (very consistent).\n"
     "PCA shows clear chromosomal clustering.\n"
     "0 out of 21,368 CpGs are missing — complete coverage."),
]):
    left = Inches(0.4 + i * 4.3)
    add_rect(slide, left, Inches(3.6), Inches(4.1), Inches(3.6), ACCENT)
    add_textbox(slide, title, left + Inches(0.15), Inches(3.7), Inches(3.8), Inches(0.45),
                font_size=15, bold=True, color=HIGHLIGHT)
    add_textbox(slide, body, left + Inches(0.15), Inches(4.2), Inches(3.8), Inches(2.8),
                font_size=13, color=LIGHT_GRAY)


# =============================================================================
# SLIDE 7 — MethylLlama Architecture
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Step 2: MethylLlama — The Model",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(5.5), Inches(0.04), HIGHLIGHT)

# Left: architecture diagram as text
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.9), ACCENT)
add_textbox(slide, "Architecture", Inches(0.6), Inches(1.3), Inches(5.4), Inches(0.4),
            font_size=18, bold=True, color=ORANGE)

arch_lines = [
    ("INPUT  [CLS, cpg₁, cpg₂, ..., cpg₁₀₆₈₄]", LIGHT_GRAY),
    ("   50% of all CpGs, randomly selected", SUBTITLE_GRAY),
    ("", WHITE),
    ("EMBEDDING LAYER", HIGHLIGHT),
    ("   h = cpg_scale × CpGEmbed(id) + ScaleAdapt(β)", LIGHT_GRAY),
    ("   ↑ 768-dim vector per CpG", SUBTITLE_GRAY),
    ("   ↑ BMFM-DNA init goes here", GREEN),
    ("", WHITE),
    ("8 × LLAMA LAYER", HIGHLIGHT),
    ("   RMSNorm + RoPE Attention + SwiGLU MLP", LIGHT_GRAY),
    ("", WHITE),
    ("CLS TOKEN  →  pooler_output  [B, 768]", ORANGE),
    ("   Global methylation representation", SUBTITLE_GRAY),
]

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.4), Inches(5.1))
tf = txBox.text_frame
tf.word_wrap = True
first = True
for text, color in arch_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = color

# Right: key design choices
add_rect(slide, Inches(6.5), Inches(1.2), Inches(6.4), Inches(5.9), ACCENT)
add_textbox(slide, "Design Choices vs Standard BERT", Inches(6.7), Inches(1.3),
            Inches(6), Inches(0.4), font_size=18, bold=True, color=ORANGE)

choices = [
    ("RoPE", "Relative position encoding inside attention heads. "
              "No position embedding table → saves parameters, handles any sequence length."),
    ("RMSNorm", "40% cheaper than LayerNorm. Normalizes by RMS only, "
                 "no mean subtraction. More stable in practice."),
    ("SwiGLU MLP", "Gated activation: SiLU(gate) × up → down. "
                    "More expressive than GELU MLP. Used in LLaMA, GPT-4."),
    ("Pre-LN", "Normalize BEFORE each sublayer, not after. "
                "Stable training from step 1 — no warmup instability."),
    ("ScaleAdapt", "Trainable sinusoidal basis for beta values [0,1]. "
                    "Special tokens (MASK, CLS) use learned embeddings via negative indices."),
]

y_off = Inches(1.85)
for title, body in choices:
    add_textbox(slide, title, Inches(6.7), y_off, Inches(1.8), Inches(0.3),
                font_size=14, bold=True, color=GREEN)
    add_textbox(slide, body, Inches(8.6), y_off, Inches(4.1), Inches(0.75),
                font_size=12, color=LIGHT_GRAY)
    y_off += Inches(1.0)


# =============================================================================
# SLIDE 8 — Injecting DNA Embeddings
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Step 3: Injecting BMFM-DNA Embeddings into MethylLlama",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=30, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(8.5), Inches(0.04), HIGHLIGHT)

# Embedding table visual
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.5), Inches(5.7), ACCENT)
add_textbox(slide, "CpG Embedding Table  [21,373 × 768]",
            Inches(0.6), Inches(1.3), Inches(5.2), Inches(0.4),
            font_size=16, bold=True, color=ORANGE)

rows_emb = [
    ("rows 0–4", "Special tokens", "UNK, SEP, PAD, CLS, MASK", RGBColor(0x80,0x80,0x80), "NEVER touched"),
    ("rows 5–21372", "CpG sites", "cg00000292 … cg27573027", GREEN, "← BMFM-DNA vectors written here"),
]

y_e = Inches(1.85)
for row_label, type_label, example, color, note in rows_emb:
    add_rect(slide, Inches(0.6), y_e, Inches(5.1), Inches(1.1),
             RGBColor(0x10, 0x20, 0x40) if "Special" in type_label else RGBColor(0x05, 0x30, 0x25))
    add_textbox(slide, row_label, Inches(0.7), y_e + Inches(0.05), Inches(1.4), Inches(0.4),
                font_size=13, bold=True, color=color)
    add_textbox(slide, type_label, Inches(2.2), y_e + Inches(0.05), Inches(1.5), Inches(0.35),
                font_size=12, color=LIGHT_GRAY)
    add_textbox(slide, example, Inches(0.7), y_e + Inches(0.5), Inches(5), Inches(0.35),
                font_size=11, color=SUBTITLE_GRAY, italic=True)
    add_textbox(slide, note, Inches(0.7), y_e + Inches(0.78), Inches(5), Inches(0.28),
                font_size=12, bold=True, color=color)
    y_e += Inches(1.25)

# Verification box
add_rect(slide, Inches(0.6), Inches(4.4), Inches(5.1), Inches(2.2), RGBColor(0x02,0x20,0x15))
add_textbox(slide, "Automatic Verification at Every Run",
            Inches(0.8), Inches(4.5), Inches(4.8), Inches(0.4),
            font_size=14, bold=True, color=GREEN)
add_textbox(slide,
    "After loading, embedding norms are checked:\n\n"
    "  DNA-initialized  →  norm ≈ 7.62  (tight)\n"
    "  Random Xavier    →  norm ≈ 0.05  (small)\n\n"
    "Logged to WandB as  dna_init/verified_ok",
    Inches(0.8), Inches(4.95), Inches(4.8), Inches(1.5),
    font_size=13, color=LIGHT_GRAY)

# Right side — explanation
add_rect(slide, Inches(6.2), Inches(1.2), Inches(6.7), Inches(5.7), ACCENT)
add_textbox(slide, "What This Gives the Model",
            Inches(6.4), Inches(1.3), Inches(6.3), Inches(0.4),
            font_size=18, bold=True, color=HIGHLIGHT)

points = [
    ("Biological similarity encoded", HIGHLIGHT2,
     "Two CpGs in the same CpG island will have similar embedding vectors from day 1 — "
     "before seeing any methylation data."),
    ("Faster convergence expected", ORANGE,
     "The model starts with a meaningful prior, rather than spending epochs "
     "learning basic CpG relationships from scratch."),
    ("Genome-scale context", GREEN,
     "Each embedding captures ±512 bp of genomic context: regulatory elements, "
     "chromatin signals, transcription factor motifs — all encoded in the DNA."),
    ("Frozen during pretraining?", HIGHLIGHT2,
     "No — embeddings are trainable. BMFM-DNA provides the starting point, "
     "but the model can refine them using methylation data."),
]

y_r = Inches(1.85)
for title, tc, body in points:
    add_textbox(slide, title, Inches(6.4), y_r, Inches(6.3), Inches(0.35),
                font_size=14, bold=True, color=tc)
    add_textbox(slide, body, Inches(6.4), y_r + Inches(0.35), Inches(6.3), Inches(0.8),
                font_size=13, color=LIGHT_GRAY)
    y_r += Inches(1.3)


# =============================================================================
# SLIDE 9 — WCED Training Method
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Step 4: Pretraining with WCED",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(5), Inches(0.04), HIGHLIGHT)
add_textbox(slide, "Weighted CpG Expression Denoising — a Masked Autoencoder for methylation",
            Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.35),
            font_size=14, color=SUBTITLE_GRAY, italic=True)

# Pipeline
steps2 = [
    ("Random 50%\nof CpGs\nas INPUT", ORANGE),
    ("MethylLlama\nEncoder\n73.7M params", HIGHLIGHT),
    ("CLS\nRepresentation\n[B, 768]", HIGHLIGHT2),
    ("Linear\nDecoder\n768 → 21368", HIGHLIGHT),
    ("Predicted\nbeta values\nfor ALL CpGs", GREEN),
    ("Loss on\nheld-out 50%\nonly", ORANGE),
]

box_w2 = Inches(1.8)
box_h2 = Inches(1.3)
gap2   = Inches(0.22)
sx2    = Inches(0.3)
y2     = Inches(1.6)

for i, (label, color) in enumerate(steps2):
    x = sx2 + i * (box_w2 + gap2)
    add_rect(slide, x, y2, box_w2, box_h2, ACCENT)
    add_rect(slide, x, y2, box_w2, Inches(0.07), color)
    add_textbox(slide, label, x + Inches(0.08), y2 + Inches(0.1),
                box_w2 - Inches(0.1), box_h2 - Inches(0.1),
                font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < len(steps2) - 1:
        ax = x + box_w2 + Inches(0.02)
        add_textbox(slide, "→", ax, y2 + Inches(0.35), gap2 + Inches(0.08), Inches(0.5),
                    font_size=18, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Why random masking
add_rect(slide, Inches(0.4), Inches(3.1), Inches(6.1), Inches(3.9), ACCENT)
add_textbox(slide, "Why Random 50% Every Batch?",
            Inches(0.6), Inches(3.2), Inches(5.8), Inches(0.4),
            font_size=17, bold=True, color=ORANGE)
add_textbox(slide,
    "If the same CpGs were always input and the same CpGs always target, "
    "the model would just memorize fixed input→output mappings.\n\n"
    "By randomly shuffling which 50% are input vs target at each batch, "
    "the model is forced to learn global methylation patterns — "
    "it can never rely on specific CpGs always being available.\n\n"
    "At inference: feed any subset of CpGs → robust representation.",
    Inches(0.6), Inches(3.65), Inches(5.8), Inches(3.1),
    font_size=14, color=LIGHT_GRAY)

# Contrastive loss
add_rect(slide, Inches(6.7), Inches(3.1), Inches(6.1), Inches(3.9), ACCENT)
add_textbox(slide, "Why Contrastive Loss?",
            Inches(6.9), Inches(3.2), Inches(5.8), Inches(0.4),
            font_size=17, bold=True, color=GREEN)
add_textbox(slide,
    "Without contrastive loss, the CLS representation collapses — "
    "the encoder learns to produce nearly identical vectors for all samples.\n\n"
    "The reconstruction task alone does not require diverse representations.\n\n"
    "InfoNCE contrastive loss: two different random 50% views of the same sample "
    "should have similar CLS embeddings, while different samples should be pushed apart.\n\n"
    "Weight = 0.1  →  Total loss = Reconstruction + 0.1 × InfoNCE",
    Inches(6.9), Inches(3.65), Inches(5.8), Inches(3.1),
    font_size=14, color=LIGHT_GRAY)


# =============================================================================
# SLIDE 10 — Full Pipeline & Experiment Design
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "Full Pipeline & Experiment Design",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(6), Inches(0.04), HIGHLIGHT)

# Pipeline flow top
pipeline = [
    "hg38.fa\n+\nManifest",
    "Extract\n±512bp\nper CpG",
    "BMFM-DNA\n113M\n(frozen)",
    "21,368\n768-dim\nvectors",
    "Init\nCpG\nEmbedding\nTable",
    "WCED\nPretraining\n300 epochs",
    "Best\nCheckpoint",
    "Fine-tune\n(age / cell type)",
]
cols = [HIGHLIGHT2, HIGHLIGHT, ORANGE, GREEN, HIGHLIGHT, HIGHLIGHT2, GREEN, ORANGE]
bw = Inches(1.45)
bh = Inches(1.2)
gp = Inches(0.15)
sx = Inches(0.3)
yp = Inches(1.2)

for i, (label, color) in enumerate(zip(pipeline, cols)):
    x = sx + i * (bw + gp)
    add_rect(slide, x, yp, bw, bh, ACCENT)
    add_rect(slide, x, yp, bw, Inches(0.06), color)
    add_textbox(slide, label, x + Inches(0.05), yp + Inches(0.08),
                bw - Inches(0.05), bh - Inches(0.08),
                font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        ax = x + bw
        add_textbox(slide, "→", ax, yp + Inches(0.35), gp + Inches(0.05), Inches(0.4),
                    font_size=14, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Two experiment boxes
add_rect(slide, Inches(0.4), Inches(2.65), Inches(6.0), Inches(4.2), ACCENT)
add_rect(slide, Inches(0.4), Inches(2.65), Inches(6.0), Inches(0.06), HIGHLIGHT2)
add_textbox(slide, "Experiment A — Baseline (Random Init)",
            Inches(0.6), Inches(2.72), Inches(5.7), Inches(0.45),
            font_size=17, bold=True, color=HIGHLIGHT2)
add_textbox(slide,
    "MethylLlama with Xavier random CpG embeddings.\n\n"
    "Identical architecture and training hyperparameters.\n\n"
    "Train: 8,724 samples\nVal: 2,264 samples\n"
    "LR: 1e-3  |  Batch: 256  |  Contrastive: 0.1\n\n"
    "Purpose: establishes baseline performance for comparison.",
    Inches(0.6), Inches(3.2), Inches(5.7), Inches(3.4),
    font_size=14, color=LIGHT_GRAY)

add_rect(slide, Inches(7.0), Inches(2.65), Inches(6.0), Inches(4.2), ACCENT)
add_rect(slide, Inches(7.0), Inches(2.65), Inches(6.0), Inches(0.06), GREEN)
add_textbox(slide, "Experiment B — BMFM-DNA Init",
            Inches(7.2), Inches(2.72), Inches(5.7), Inches(0.45),
            font_size=17, bold=True, color=GREEN)
add_textbox(slide,
    "Same MethylLlama, but CpG embedding table rows 5-21372\nare initialized with BMFM-DNA vectors.\n\n"
    "Identical architecture and training hyperparameters.\n\n"
    "Train: 8,724 samples\nVal: 2,264 samples\n"
    "LR: 1e-3  |  Batch: 256  |  Contrastive: 0.1\n\n"
    "Purpose: test whether DNA context gives better representations.",
    Inches(7.2), Inches(3.2), Inches(5.7), Inches(3.4),
    font_size=14, color=LIGHT_GRAY)

# VS in middle
add_textbox(slide, "VS", Inches(6.05), Inches(4.1), Inches(0.9), Inches(0.7),
            font_size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 11 — What We Expect to See
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), HIGHLIGHT)

add_textbox(slide, "What We Expect to See",
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=HIGHLIGHT)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(4.5), Inches(0.04), HIGHLIGHT)

metrics = [
    ("Faster Convergence", ORANGE,
     "DNA-init should reach low validation loss in fewer epochs — "
     "the model starts with biologically meaningful CpG similarities rather than learning them from scratch."),
    ("Better Reconstruction (val/pcc)", HIGHLIGHT2,
     "Higher Pearson correlation between predicted and true beta values on held-out CpGs. "
     "Current run: val/pcc ≈ 0.96 at epoch 5 — expect DNA-init to go higher."),
    ("More Discriminative Representations", GREEN,
     "Higher cls_variance in WandB — the CLS embeddings should be more spread out, "
     "capturing cell type, tissue, and age differences more clearly than random init."),
    ("Better Downstream Fine-tuning", ORANGE,
     "The ultimate test: fine-tune both checkpoints on biological age prediction. "
     "If DNA-init wins, the genomic context prior is genuinely useful for methylation modeling."),
]

y_m = Inches(1.2)
for title, color, body in metrics:
    add_rect(slide, Inches(0.4), y_m, Inches(12.5), Inches(1.35), ACCENT)
    add_rect(slide, Inches(0.4), y_m, Inches(0.08), Inches(1.35), color)
    add_textbox(slide, title, Inches(0.65), y_m + Inches(0.1), Inches(4), Inches(0.45),
                font_size=16, bold=True, color=color)
    add_textbox(slide, body, Inches(0.65), y_m + Inches(0.55), Inches(12), Inches(0.7),
                font_size=14, color=LIGHT_GRAY)
    y_m += Inches(1.5)

add_textbox(slide,
    "Current status (epoch 20):  val/pcc = 0.958  •  val/loss = 0.00874  •  Training running ✓",
    Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
    font_size=13, color=GREEN, bold=True)


# =============================================================================
# Save
# =============================================================================
out_path = "/Users/netanelazran/Projects/BMFM-RNA_thesis/methyl/BMFM_DNA_MethylLlama_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
