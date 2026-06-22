"""
Generate PowerPoint: Data Pipeline & NaN Handling in MethylLlama
================================================================
White-background, diagram-first version.
No raw code panels — every step shown as visual flow / table / diagram.

Usage:
    python scripts/utils/make_nan_presentation.py
    # → NaN_Handling_MethylLlama.pptx
"""

import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── White-style palette ────────────────────────────────────────────────────────
BG           = RGBColor(0xFF, 0xFF, 0xFF)   # pure white
PANEL_BG     = RGBColor(0xF2, 0xF6, 0xFB)  # very light blue-gray
PANEL_BG2    = RGBColor(0xE8, 0xF4, 0xED)  # very light green
PANEL_BG3    = RGBColor(0xFD, 0xF3, 0xE4)  # very light orange
PANEL_BG4    = RGBColor(0xFD, 0xEC, 0xEC)  # very light red
TEXT_DARK    = RGBColor(0x1A, 0x26, 0x3A)  # dark navy — headings
TEXT_BODY    = RGBColor(0x3A, 0x4A, 0x5C)  # medium navy — body
TEXT_LIGHT   = RGBColor(0x78, 0x8F, 0xA8)  # muted — captions
DIVIDER      = RGBColor(0xCC, 0xD9, 0xE8)  # light border
BLUE         = RGBColor(0x00, 0x7A, 0xC2)  # primary accent
BLUE_DARK    = RGBColor(0x00, 0x4F, 0x80)  # darker blue
GREEN        = RGBColor(0x00, 0x8F, 0x6A)  # green
GREEN_DARK   = RGBColor(0x00, 0x60, 0x45)  # darker green
ORANGE       = RGBColor(0xD4, 0x7E, 0x00)  # amber
ORANGE_DARK  = RGBColor(0x9A, 0x5A, 0x00)
RED          = RGBColor(0xC0, 0x30, 0x30)  # red
PURPLE       = RGBColor(0x6A, 0x28, 0x9A)  # purple
BADGE_BLUE   = RGBColor(0xD6, 0xEA, 0xF8)  # pastel blue for badge bg
BADGE_GREEN  = RGBColor(0xD5, 0xF0, 0xE4)  # pastel green
BADGE_ORANGE = RGBColor(0xFD, 0xEA, 0xC8)  # pastel orange
BADGE_RED    = RGBColor(0xFA, 0xD7, 0xD7)   # pastel red

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def slide_header(slide, title, subtitle=None):
    """Top bar + title + optional subtitle."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), BLUE)
    add_textbox(slide, title,
                Inches(0.55), Inches(0.15), Inches(12.2), Inches(0.65),
                font_size=28, bold=True, color=TEXT_DARK)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.55), Inches(0.82), Inches(12.2), Inches(0.35),
                    font_size=14, color=TEXT_LIGHT, italic=True)
    # bottom divider
    add_rect(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.025), DIVIDER)


def card(slide, left, top, width, height, bg=PANEL_BG,
         border_color=None, radius=False):
    """Flat card (rounded not supported in python-pptx, just a rect)."""
    return add_rect(slide, left, top, width, height, bg,
                    line_color=border_color or DIVIDER)


def badge(slide, text, left, top, width, height,
          bg=BADGE_BLUE, fg=BLUE_DARK, font_size=13, bold=False):
    add_rect(slide, left, top, width, height, bg, line_color=None)
    add_textbox(slide, text, left, top, width, height,
                font_size=font_size, bold=bold, color=fg, align=PP_ALIGN.CENTER)


def section_label(slide, text, left, top, color=BLUE):
    """Small ALL-CAPS section label."""
    add_textbox(slide, text, left, top, Inches(6), Inches(0.3),
                font_size=11, bold=True, color=color)


def fig_to_slide(slide, fig, left, top, width, height):
    """Embed a matplotlib figure into a slide."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)


def arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.4)):
    add_textbox(slide, "→", left, top, width, height,
                font_size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def arrow_down(slide, left, top, width=Inches(0.5), height=Inches(0.4)):
    add_textbox(slide, "▼", left, top, width, height,
                font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def icon_box(slide, icon, label, left, top, width, height,
             icon_color=BLUE, bg=PANEL_BG):
    card(slide, left, top, width, height, bg=bg)
    add_textbox(slide, icon, left, top + Inches(0.1), width, Inches(0.55),
                font_size=28, color=icon_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + Inches(0.62), width, Inches(0.45),
                font_size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def metric_card(slide, left, top, width, height, value, label,
                val_color=BLUE, bg=PANEL_BG):
    card(slide, left, top, width, height, bg=bg, border_color=val_color)
    add_textbox(slide, value, left, top + Inches(0.08), width, Inches(0.55),
                font_size=30, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + Inches(0.58), width, Inches(0.45),
                font_size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)


def flow_box(slide, left, top, width, height, title, body,
             accent=BLUE, bg=PANEL_BG):
    card(slide, left, top, width, height, bg=bg, border_color=None)
    add_rect(slide, left, top, Inches(0.06), height, accent)
    add_textbox(slide, title, left + Inches(0.14), top + Inches(0.1),
                width - Inches(0.2), Inches(0.38),
                font_size=13, bold=True, color=accent)
    add_textbox(slide, body, left + Inches(0.14), top + Inches(0.48),
                width - Inches(0.2), height - Inches(0.55),
                font_size=12, color=TEXT_BODY)


def table_row(slide, cols, y, col_widths, col_lefts,
              fg=TEXT_BODY, bg=None, font_size=13, bold=False):
    if bg:
        add_rect(slide, col_lefts[0], y, sum(col_widths), Inches(0.42), bg)
    for text, w, x in zip(cols, col_widths, col_lefts):
        add_textbox(slide, text, x + Inches(0.1), y + Inches(0.05),
                    w - Inches(0.2), Inches(0.35),
                    font_size=font_size, color=fg, bold=bold)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), BLUE)

add_textbox(slide, "Data Pipeline & NaN Handling",
            Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.1),
            font_size=46, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, "MethylLlama Pretraining on the 49k CpG Union Corpus",
            Inches(1.5), Inches(2.75), Inches(10.3), Inches(0.6),
            font_size=22, color=BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.0), Inches(3.55), Inches(5.3), Inches(0.04), DIVIDER)
add_textbox(slide,
    "How missing measurements (NaN) flow through the pipeline\n"
    "and are correctly excluded from the reconstruction loss",
    Inches(2.0), Inches(3.72), Inches(9.3), Inches(0.7),
    font_size=17, color=TEXT_BODY, align=PP_ALIGN.CENTER)

for i, (line, col, sz) in enumerate([
    ("Netanel Azran  |  BMFM-RNA Methylation  |  Hebrew University", TEXT_LIGHT, 14),
    ("Run: llama-small-all49k-r0.5-w0.0-44450919", BLUE, 13),
]):
    add_textbox(slide, line, Inches(1.5), Inches(5.6 + i * 0.38),
                Inches(10.3), Inches(0.38), font_size=sz, color=col, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — The Two Datasets
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "The Two Datasets",
             "Pretrain: 169k samples × 49k CpGs   |   Fine-tune: 11.5k samples × 49k CpGs")

# Left — pretrain card
card(slide, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.7),
     bg=PANEL_BG, border_color=BLUE)
add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.9), Inches(0.06), BLUE)
badge(slide, "PRETRAINING", Inches(0.55), Inches(1.5), Inches(1.5), Inches(0.3),
      bg=BADGE_BLUE, fg=BLUE_DARK, font_size=11, bold=True)
add_textbox(slide, "169,120 samples × 49,156 CpGs",
            Inches(0.6), Inches(1.92), Inches(5.5), Inches(0.45),
            font_size=20, bold=True, color=TEXT_DARK)

rows_pre = [
    ("📁 File",       "methylgpt_pretrain_type3.h5ad"),
    ("⚙  Script",    "pretrain_llama_small.sh  →  pretrain_llama.py"),
    ("🔷 Module",     "WCEDLlamaModule  (wced_llama.py)"),
    ("📐 Architecture", "256D × 4L × 4H  ≈ 5M params"),
    ("🎯 Input ratio", "0.5  — 50% of valid CpGs per view"),
    ("🏋  Age weight", "0.0  — pure reconstruction pretraining"),
    ("🔖 Job ID",     "44450919  (7-day run)"),
]
for i, (k, v) in enumerate(rows_pre):
    y = Inches(2.52 + i * 0.46)
    add_textbox(slide, k, Inches(0.65), y, Inches(1.6), Inches(0.4),
                font_size=12, bold=True, color=BLUE)
    add_textbox(slide, v, Inches(2.3), y, Inches(3.85), Inches(0.4),
                font_size=12, color=TEXT_BODY)

# Right — finetune card
card(slide, Inches(7.0), Inches(1.4), Inches(5.9), Inches(5.7),
     bg=PANEL_BG2, border_color=GREEN)
add_rect(slide, Inches(7.0), Inches(1.4), Inches(5.9), Inches(0.06), GREEN)
badge(slide, "FINE-TUNING", Inches(7.15), Inches(1.5), Inches(1.5), Inches(0.3),
      bg=BADGE_GREEN, fg=GREEN_DARK, font_size=11, bold=True)
add_textbox(slide, "11,500 samples × 49,156 CpGs",
            Inches(7.2), Inches(1.92), Inches(5.5), Inches(0.45),
            font_size=20, bold=True, color=TEXT_DARK)

rows_ft = [
    ("📁 File",       "finetuning_49k.h5ad"),
    ("⚙  Script",    "finetune_llama_small.sh  →  finetune_llama.py"),
    ("🔷 Module",     "WCEDLlamaModule  (wced_llama.py)"),
    ("📐 Pooling",    "Mean over all tokens"),
    ("🎯 Input ratio", "1.0  — all valid CpGs as input"),
    ("🏋  Recon weight", "0.0  — age MSE loss only"),
    ("❄  Encoder",   "Frozen  (unfrozen after epoch 10)"),
]
for i, (k, v) in enumerate(rows_ft):
    y = Inches(2.52 + i * 0.46)
    add_textbox(slide, k, Inches(7.25), y, Inches(1.6), Inches(0.4),
                font_size=12, bold=True, color=GREEN)
    add_textbox(slide, v, Inches(8.9), y, Inches(3.85), Inches(0.4),
                font_size=12, color=TEXT_BODY)

# Middle connector
add_textbox(slide, "Same\n49k vocab", Inches(6.08), Inches(3.6), Inches(0.85), Inches(0.75),
            font_size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 3 — Why NaN Exists
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Why NaN Exists — The Illumina Array Union Problem",
             "Each sample was measured on ONE array type → other probes are structurally missing")

# Three array cards
for i, (array, probes, year, color, bg) in enumerate([
    ("27k Array",   "~27,000 probes", "Legacy",   ORANGE,  BADGE_ORANGE),
    ("450k Array",  "~450,000 probes","2012–",     BLUE,    BADGE_BLUE),
    ("EPIC Array",  "~850,000 probes","2016–",     GREEN,   BADGE_GREEN),
]):
    lx = Inches(0.5 + i * 3.0)
    card(slide, lx, Inches(1.45), Inches(2.7), Inches(1.55), bg=bg, border_color=color)
    add_textbox(slide, array, lx + Inches(0.15), Inches(1.55), Inches(2.4), Inches(0.42),
                font_size=18, bold=True, color=color)
    add_textbox(slide, probes, lx + Inches(0.15), Inches(1.98), Inches(2.4), Inches(0.35),
                font_size=14, color=TEXT_BODY)
    add_textbox(slide, year, lx + Inches(0.15), Inches(2.3), Inches(2.4), Inches(0.3),
                font_size=12, color=TEXT_LIGHT)
    arrow_down(slide, lx + Inches(1.1), Inches(3.03))

# Union box
card(slide, Inches(0.5), Inches(3.46), Inches(8.7), Inches(0.75),
     bg=BADGE_BLUE, border_color=BLUE)
add_textbox(slide, "Union of all studies  →  49,156 CpGs in shared vocabulary",
            Inches(0.65), Inches(3.56), Inches(8.3), Inches(0.5),
            font_size=18, bold=True, color=BLUE_DARK)
arrow_right(slide, Inches(9.35), Inches(3.65))
card(slide, Inches(9.9), Inches(3.46), Inches(3.0), Inches(0.75),
     bg=BADGE_RED, border_color=RED)
add_textbox(slide, "NaN where not measured",
            Inches(10.0), Inches(3.56), Inches(2.8), Inches(0.5),
            font_size=14, bold=True, color=RED)

# Per-sample breakdown table
section_label(slide, "PER-SAMPLE BREAKDOWN", Inches(0.5), Inches(4.42))
card(slide, Inches(0.5), Inches(4.75), Inches(12.4), Inches(2.35), bg=PANEL_BG)

# Table header
col_w = [Inches(3.0), Inches(2.5), Inches(2.5), Inches(2.5), Inches(1.9)]
col_x = [Inches(0.5), Inches(3.5), Inches(6.0), Inches(8.5), Inches(11.0)]
table_row(slide,
    ["Sample type", "Valid CpGs", "NaN positions", "NaN rate", "Typical in"],
    Inches(4.75), col_w, col_x,
    fg=TEXT_LIGHT, font_size=12, bold=True)
add_rect(slide, Inches(0.5), Inches(5.18), Inches(12.4), Inches(0.02), DIVIDER)

rows = [
    ("450k-array sample", "~19,600", "~29,556", "~60%", RED,   "Fine-tune dataset (confirmed)"),
    ("EPIC-array sample",  "~40,000", "~9,156",  "~19%", ORANGE,"Most of pretrain corpus"),
    ("27k-array sample",   "~27,000", "~22,156", "~45%", ORANGE,"Some of pretrain corpus"),
]
for i, (samp, valid, nan, rate, rc, tip) in enumerate(rows):
    y = Inches(5.22 + i * 0.58)
    bg2 = PANEL_BG if i % 2 == 0 else BG
    add_rect(slide, Inches(0.5), y, Inches(12.4), Inches(0.56), bg2)
    for text, w, x, fc in [
        (samp, col_w[0], col_x[0], TEXT_BODY),
        (valid, col_w[1], col_x[1], GREEN),
        (nan,   col_w[2], col_x[2], RED),
        (rate,  col_w[3], col_x[3], rc),
        (tip,   col_w[4], col_x[4], TEXT_LIGHT),
    ]:
        add_textbox(slide, text, x + Inches(0.1), y + Inches(0.1),
                    w - Inches(0.15), Inches(0.38), font_size=13, color=fc)

add_textbox(slide,
    "★  From pretrain valid_pct ≈ 42%  →  NaN rate ≈ 16%  "
    "(pretrain corpus is mostly well-measured EPIC-array samples)",
    Inches(0.5), Inches(7.1), Inches(12.4), Inches(0.32),
    font_size=13, italic=True, color=GREEN)


# =============================================================================
# SLIDE 4 — Full Pipeline Overview
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Full Pipeline — NaN Handling at Each Stage",
             "From raw file to loss function: where NaN is detected, imputed, and excluded")

steps = [
    ("h5ad File",           "methylgpt_pretrain\n_type3.h5ad\n169k × 49k",   BLUE,   "Input",   BADGE_BLUE),
    ("MethylationDataset",  "Reads β values\nComputes valid_mask\n= isfinite(β)",ORANGE,"Step 1", BADGE_ORANGE),
    ("WCEDCollator  2a",    "NaN → 0 in targets\nvalid_mask stored\nin batch dict",GREEN,"Step 2a",BADGE_GREEN),
    ("WCEDCollator  2b",    "Input views built\nfrom valid CpGs only\n(NaN never tokenised)",PURPLE,"Step 2b",RGBColor(0xEE,0xE0,0xF8)),
    ("WCEDLlamaModule",     "recon_mask =\n~input_mask\n& valid_mask",        BLUE,   "Step 3",  BADGE_BLUE),
]

bw = Inches(2.22)
bh = Inches(1.9)
gap = Inches(0.18)
sx = Inches(0.3)
sy = Inches(1.45)

for i, (title, body, color, tag, bg) in enumerate(steps):
    lx = sx + i * (bw + gap)
    card(slide, lx, sy, bw, bh, bg=bg, border_color=color)
    badge(slide, tag, lx + Inches(0.1), sy + Inches(0.08),
          Inches(0.9), Inches(0.27), bg=color, fg=BG, font_size=10, bold=True)
    add_textbox(slide, title, lx + Inches(0.1), sy + Inches(0.38),
                bw - Inches(0.2), Inches(0.4), font_size=13, bold=True, color=color)
    add_textbox(slide, body, lx + Inches(0.1), sy + Inches(0.82),
                bw - Inches(0.2), bh - Inches(0.9), font_size=11, color=TEXT_BODY)
    if i < len(steps) - 1:
        arrow_right(slide, lx + bw + Inches(0.01), sy + Inches(0.72))

# Three outcome lanes
section_label(slide, "CpG POSITIONS — THREE CATEGORIES PER SAMPLE", Inches(0.4), Inches(3.58))

for i, (label, desc, outcome, color, bg) in enumerate([
    ("In input view",    "~20.5k positions\n50% of valid CpGs\nshown to encoder",
     "EXCLUDED from recon loss\n(encoder already saw it)", ORANGE, BADGE_ORANGE),
    ("Valid — held out", "~20.5k positions\n50% of valid CpGs\nnot in input",
     "INCLUDED in recon loss ✓\nReal β target  →  true signal", GREEN, BADGE_GREEN),
    ("NaN position",     "~8k positions\nnever measured\ntarget = 0.0 (imputed)",
     "EXCLUDED by valid_mask ✓\nNo gradient from fake zeros", RED, BADGE_RED),
]):
    lx = Inches(0.4 + i * 4.3)
    card(slide, lx, Inches(3.95), Inches(4.1), Inches(3.2), bg=bg, border_color=color)
    add_rect(slide, lx, Inches(3.95), Inches(4.1), Inches(0.05), color)
    add_textbox(slide, label, lx + Inches(0.2), Inches(4.05),
                Inches(3.7), Inches(0.4), font_size=15, bold=True, color=color)
    add_textbox(slide, desc, lx + Inches(0.2), Inches(4.52),
                Inches(3.7), Inches(0.88), font_size=12, color=TEXT_BODY)
    add_rect(slide, lx + Inches(0.2), Inches(5.44), Inches(3.4), Inches(0.025), color)
    add_textbox(slide, outcome, lx + Inches(0.2), Inches(5.5),
                Inches(3.7), Inches(0.6), font_size=12, bold=True, color=color)


# =============================================================================
# SLIDE 5 — Step 1: MethylationDataset
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Step 1 — MethylationDataset: Reading the Data",
             "bmfm_methylation/shared/data_module.py  ·  MethylationDataset.__getitem__")

# Flow diagram
flow_nodes = [
    ("📄  adata.X[idx]", "Raw row from h5ad\n49,156 float values\nSparse or dense"),
    ("🔢  float32 array", "Convert to numpy\nfloat32 array\nshape: (49156,)"),
    ("🔍  isfinite(β)", "np.isfinite(beta_values)\nTrue  = real value\nFalse = NaN"),
]
node_w = Inches(2.8)
node_h = Inches(1.8)
for i, (title, body) in enumerate(flow_nodes):
    lx = Inches(0.4 + i * 3.3)
    card(slide, lx, Inches(1.45), node_w, node_h, bg=PANEL_BG, border_color=BLUE)
    add_textbox(slide, title, lx + Inches(0.15), Inches(1.55), node_w - Inches(0.3),
                Inches(0.45), font_size=13, bold=True, color=BLUE)
    add_textbox(slide, body, lx + Inches(0.15), Inches(2.02), node_w - Inches(0.3),
                Inches(1.1), font_size=12, color=TEXT_BODY)
    if i < len(flow_nodes) - 1:
        arrow_right(slide, Inches(0.4 + i * 3.3) + node_w + Inches(0.05),
                    Inches(2.1))

# Split into two outputs
arrow_down(slide, Inches(3.78), Inches(3.28))
add_textbox(slide, "Two outputs", Inches(3.1), Inches(3.72), Inches(2.15), Inches(0.3),
            font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, "↙         ↘", Inches(3.1), Inches(3.98), Inches(2.15), Inches(0.4),
            font_size=20, color=BLUE, align=PP_ALIGN.CENTER)

# Left output: beta_values
card(slide, Inches(2.0), Inches(4.42), Inches(2.8), Inches(1.65),
     bg=BADGE_ORANGE, border_color=ORANGE)
add_textbox(slide, "beta_values", Inches(2.15), Inches(4.52), Inches(2.5),
            Inches(0.38), font_size=13, bold=True, color=ORANGE)
add_textbox(slide,
    "Shape: (49156,)  float32\n"
    "NaN values preserved\n"
    "Passed to collator as-is",
    Inches(2.15), Inches(4.92), Inches(2.5), Inches(1.1),
    font_size=12, color=TEXT_BODY)

# Right output: valid_mask
card(slide, Inches(5.4), Inches(4.42), Inches(2.8), Inches(1.65),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide, "valid_mask", Inches(5.55), Inches(4.52), Inches(2.5),
            Inches(0.38), font_size=13, bold=True, color=GREEN)
add_textbox(slide,
    "Shape: (49156,)  bool\n"
    "True = measured\n"
    "False = NaN position",
    Inches(5.55), Inches(4.92), Inches(2.5), Inches(1.1),
    font_size=12, color=TEXT_BODY)

# Right: why card
card(slide, Inches(9.2), Inches(1.45), Inches(3.9), Inches(4.65),
     bg=PANEL_BG, border_color=DIVIDER)
section_label(slide, "WHY THIS DESIGN", Inches(9.35), Inches(1.58), color=ORANGE)
for i, point in enumerate([
    "NaN is preserved — not replaced here",
    "valid_mask is computed while NaN is still visible",
    "If we replaced NaN→0 first, the mask would be wrong",
    "Separation of concerns: Dataset reads, Collator decides",
    "valid_mask travels through the batch dict to the loss fn",
]):
    add_textbox(slide, f"{'●'}  {point}",
                Inches(9.35), Inches(1.98 + i * 0.74), Inches(3.55), Inches(0.65),
                font_size=13, color=TEXT_BODY)

# Bottom status
card(slide, Inches(0.4), Inches(6.28), Inches(8.65), Inches(0.55),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide,
    "After Step 1:   beta_values contains NaN  ·  valid_mask = True/False array  ·  both in MultiFieldInstance",
    Inches(0.55), Inches(6.35), Inches(8.35), Inches(0.4),
    font_size=13, bold=True, color=GREEN_DARK)


# =============================================================================
# SLIDE 6 — Step 2a: WCEDCollator — NaN Detection & Replacement
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Step 2a — WCEDCollator: NaN Detection & Replacement",
             "bmfm_methylation/shared/data_module.py  ·  WCEDCollator.__call__")

# Visual: before → after tensor
section_label(slide, "VISUAL EXAMPLE — ONE SAMPLE ROW (6 CpGs shown)", Inches(0.4), Inches(1.4))

# Before row
add_textbox(slide, "Input (beta_values):", Inches(0.4), Inches(1.72),
            Inches(2.5), Inches(0.38), font_size=13, bold=True, color=TEXT_BODY)
cells_before = [("0.23", BADGE_GREEN, GREEN_DARK), ("NaN", BADGE_RED, RED),
                ("0.87", BADGE_GREEN, GREEN_DARK), ("NaN", BADGE_RED, RED),
                ("0.61", BADGE_GREEN, GREEN_DARK), ("NaN", BADGE_RED, RED)]
for j, (val, bg, fg) in enumerate(cells_before):
    lx = Inches(3.1 + j * 1.12)
    card(slide, lx, Inches(1.72), Inches(1.0), Inches(0.38), bg=bg, border_color=None)
    add_textbox(slide, val, lx, Inches(1.72), Inches(1.0), Inches(0.38),
                font_size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)

arrow_down(slide, Inches(6.4), Inches(2.14))
add_textbox(slide, "np.where(valid, β, 0.0)", Inches(5.4), Inches(2.52),
            Inches(2.5), Inches(0.3), font_size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
arrow_down(slide, Inches(6.4), Inches(2.78))

# After: all_betas
add_textbox(slide, "all_betas (target):", Inches(0.4), Inches(3.14),
            Inches(2.5), Inches(0.38), font_size=13, bold=True, color=TEXT_BODY)
cells_after = [("0.23", BADGE_GREEN, GREEN_DARK), ("0.00", BADGE_ORANGE, ORANGE_DARK),
               ("0.87", BADGE_GREEN, GREEN_DARK), ("0.00", BADGE_ORANGE, ORANGE_DARK),
               ("0.61", BADGE_GREEN, GREEN_DARK), ("0.00", BADGE_ORANGE, ORANGE_DARK)]
for j, (val, bg, fg) in enumerate(cells_after):
    lx = Inches(3.1 + j * 1.12)
    card(slide, lx, Inches(3.14), Inches(1.0), Inches(0.38), bg=bg, border_color=None)
    add_textbox(slide, val, lx, Inches(3.14), Inches(1.0), Inches(0.38),
                font_size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)
add_textbox(slide, "← NaN replaced with 0.0", Inches(9.85), Inches(3.14),
            Inches(3.0), Inches(0.38), font_size=12, italic=True, color=ORANGE)

# valid_mask
add_textbox(slide, "valid_mask:", Inches(0.4), Inches(3.68),
            Inches(2.5), Inches(0.38), font_size=13, bold=True, color=TEXT_BODY)
cells_vm = [("True", BADGE_GREEN, GREEN_DARK), ("False", BADGE_RED, RED),
            ("True", BADGE_GREEN, GREEN_DARK), ("False", BADGE_RED, RED),
            ("True", BADGE_GREEN, GREEN_DARK), ("False", BADGE_RED, RED)]
for j, (val, bg, fg) in enumerate(cells_vm):
    lx = Inches(3.1 + j * 1.12)
    card(slide, lx, Inches(3.68), Inches(1.0), Inches(0.38), bg=bg, border_color=None)
    add_textbox(slide, val, lx, Inches(3.68), Inches(1.0), Inches(0.38),
                font_size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
add_textbox(slide, "← passed in batch dict →  Step 3", Inches(9.85), Inches(3.68),
            Inches(3.0), Inches(0.38), font_size=12, italic=True, color=GREEN)

# Two result cards
card(slide, Inches(0.4), Inches(4.28), Inches(5.85), Inches(1.9),
     bg=BADGE_ORANGE, border_color=ORANGE)
add_textbox(slide, "all_betas  — reconstruction target",
            Inches(0.55), Inches(4.35), Inches(5.5), Inches(0.38),
            font_size=14, bold=True, color=ORANGE_DARK)
add_textbox(slide,
    "Shape: [batch, 49156]   dtype: float32\n"
    "Real positions: actual β ∈ [0, 1]\n"
    "NaN positions: 0.0  (imputed — safe because valid_mask will exclude them from loss)\n"
    "This is what the model tries to predict",
    Inches(0.55), Inches(4.75), Inches(5.5), Inches(1.3),
    font_size=13, color=TEXT_BODY)

card(slide, Inches(6.7), Inches(4.28), Inches(6.25), Inches(1.9),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide, "valid_mask  — the critical guard",
            Inches(6.85), Inches(4.35), Inches(5.9), Inches(0.38),
            font_size=14, bold=True, color=GREEN_DARK)
add_textbox(slide,
    "Shape: [batch, 49156]   dtype: bool\n"
    "True  = position was actually measured  →  include in loss\n"
    "False = position was NaN                →  exclude from loss\n"
    "Returned in batch dict  →  used in Step 3 to build recon_mask",
    Inches(6.85), Inches(4.75), Inches(5.9), Inches(1.3),
    font_size=13, color=TEXT_BODY)

card(slide, Inches(0.4), Inches(6.32), Inches(12.55), Inches(0.55),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide,
    "After Step 2a:   all_betas ready (NaN→0.0)  ·  valid_mask ready (True/False)  ·  both in batch dict",
    Inches(0.55), Inches(6.39), Inches(12.2), Inches(0.4),
    font_size=13, bold=True, color=GREEN_DARK)


# =============================================================================
# SLIDE 7 — Step 2b: WCEDCollator — Input View Selection
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Step 2b — WCEDCollator: Input View Selection",
             "NaN positions are never tokenised — input is sampled from valid CpGs only")

# Segmented bar — 49k CpG vocabulary
section_label(slide, "49,156 CpG VOCABULARY — BREAKDOWN PER SAMPLE", Inches(0.4), Inches(1.38))

bar_left = Inches(0.55)
bar_top  = Inches(1.72)
bar_h    = Inches(0.72)
bar_total_w = Inches(12.2)

segments = [
    ("NaN positions\n~4.3k  (8.7%)",   0.087, BADGE_RED,    RED),
    ("Input view\n~22.5k  (45.7%)",    0.457, BADGE_ORANGE, ORANGE),
    ("Held-out (recon target)\n~22.5k  (45.7%)", 0.456, BADGE_GREEN, GREEN),
]
x = bar_left
for label, frac, bg, fg in segments:
    w = bar_total_w * frac
    add_rect(slide, x, bar_top, w, bar_h, bg)
    lw = Pt(1.5)
    shape = slide.shapes[-1]
    shape.line.color.rgb = fg
    shape.line.width = lw
    add_textbox(slide, label, x + Inches(0.05), bar_top + Inches(0.08),
                w - Inches(0.1), bar_h - Inches(0.12),
                font_size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
    x += w

# Formula
section_label(slide, "INPUT RATIO FORMULA", Inches(0.4), Inches(2.66))
card(slide, Inches(0.4), Inches(2.96), Inches(7.2), Inches(1.5),
     bg=PANEL_BG, border_color=BLUE)
add_textbox(slide,
    "valid_indices  =  positions where isfinite(β)  →  ~44.9k of 49k  (NaN rate = 8.7%)",
    Inches(0.55), Inches(3.02), Inches(6.9), Inches(0.4),
    font_size=14, color=TEXT_BODY)
add_textbox(slide,
    "n_input  =  int( len(valid_indices)  ×  input_ratio )\n"
    "         =  int( 44,900  ×  0.5 )  =  22,450",
    Inches(0.55), Inches(3.42), Inches(6.9), Inches(0.55),
    font_size=14, bold=True, color=BLUE)
add_textbox(slide,
    "input_ratio applied to valid CpGs  —  NOT to the full 49k vocabulary",
    Inches(0.55), Inches(3.98), Inches(6.9), Inches(0.38),
    font_size=12, italic=True, color=ORANGE)

# Why it matters card
card(slide, Inches(7.75), Inches(2.66), Inches(5.25), Inches(1.8),
     bg=BADGE_ORANGE, border_color=ORANGE)
add_textbox(slide, "Why this matters for high-NaN samples",
            Inches(7.9), Inches(2.73), Inches(5.0), Inches(0.38),
            font_size=14, bold=True, color=ORANGE_DARK)
add_textbox(slide,
    "450k sample (60% NaN)  →  only 19.6k valid CpGs\n"
    "n_input = 0.5 × 19.6k = 9.8k  (not 0.5 × 49k = 24.5k)\n\n"
    "Always exactly 50% of what was measured is shown.\n"
    "NaN positions NEVER appear as encoder input tokens.",
    Inches(7.9), Inches(3.12), Inches(5.0), Inches(1.25),
    font_size=13, color=TEXT_BODY)

# 3 outcome cards
section_label(slide, "WHAT GETS ASSIGNED TO EACH CATEGORY", Inches(0.4), Inches(4.7))
for i, (title, body, action, color, bg) in enumerate([
    ("In input view\n~n_input positions",
     "randomly sampled from valid_indices\ninput_ratio × |valid_indices|",
     "→ encoder tokens\n→ excluded from recon loss",
     ORANGE, BADGE_ORANGE),
    ("Valid — held out\nremaining valid positions",
     "|valid_indices| - n_input positions\nnot sampled into input view",
     "→ recon_mask = True\n→ included in loss ✓",
     GREEN, BADGE_GREEN),
    ("NaN positions\nnever in valid_indices",
     "np.where(vocab_valid) returns only\nnon-NaN indices — NaN excluded",
     "→ valid_mask = False\n→ excluded from loss ✓",
     RED, BADGE_RED),
]):
    lx = Inches(0.4 + i * 4.3)
    card(slide, lx, Inches(5.02), Inches(4.1), Inches(2.05), bg=bg, border_color=color)
    add_textbox(slide, title, lx + Inches(0.15), Inches(5.1), Inches(3.8),
                Inches(0.5), font_size=13, bold=True, color=color)
    add_textbox(slide, body, lx + Inches(0.15), Inches(5.62), Inches(3.8),
                Inches(0.5), font_size=12, color=TEXT_BODY)
    add_rect(slide, lx + Inches(0.15), Inches(6.16), Inches(3.4), Inches(0.02), color)
    add_textbox(slide, action, lx + Inches(0.15), Inches(6.2), Inches(3.8),
                Inches(0.55), font_size=12, bold=True, color=color)


# =============================================================================
# SLIDE 8 — Step 3: WCEDLlamaModule — Loss Masking
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Step 3 — WCEDLlamaModule: Correct Reconstruction Loss",
             "bmfm_methylation/llama/wced_llama.py  ·  _shared_step")

# recon_mask formula
section_label(slide, "HOW recon_mask IS BUILT", Inches(0.4), Inches(1.38))

mask_parts = [
    ("~input_mask", "Positions NOT\nshown to encoder", ORANGE, BADGE_ORANGE),
    ("AND", "", TEXT_LIGHT, BG),
    ("valid_mask", "Positions that were\nactually measured", GREEN, BADGE_GREEN),
    ("=", "", TEXT_LIGHT, BG),
    ("recon_mask", "Positions included\nin loss computation", BLUE, BADGE_BLUE),
]
x = Inches(0.5)
for title, sub, color, bg in mask_parts:
    if title in ("AND", "="):
        add_textbox(slide, title, x, Inches(1.78), Inches(0.65), Inches(0.55),
                    font_size=24, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        x += Inches(0.65)
    else:
        card(slide, x, Inches(1.72), Inches(2.3), Inches(0.85), bg=bg, border_color=color)
        add_textbox(slide, title, x + Inches(0.1), Inches(1.78), Inches(2.1),
                    Inches(0.38), font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, x + Inches(0.1), Inches(2.18), Inches(2.1),
                    Inches(0.35), font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        x += Inches(2.35)

# Decision table
section_label(slide, "POSITION-BY-POSITION DECISION TABLE", Inches(0.4), Inches(2.78))

t_left   = Inches(0.4)
t_top    = Inches(3.08)
t_w      = Inches(12.55)
t_row_h  = Inches(0.58)
col_ws   = [Inches(2.6), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.0), Inches(2.3)]
col_xs   = [Inches(0.4), Inches(3.0), Inches(4.8), Inches(6.6), Inches(8.6), Inches(10.6)]

# Header
add_rect(slide, t_left, t_top, t_w, t_row_h, PANEL_BG)
for text, w, x in zip(
    ["Position type", "~Count", "Encoder sees?", "in ~input_mask?", "valid_mask?", "In recon loss?"],
    col_ws, col_xs
):
    add_textbox(slide, text, x + Inches(0.08), t_top + Inches(0.12),
                w - Inches(0.12), Inches(0.38), font_size=13, bold=True,
                color=TEXT_DARK)

add_rect(slide, t_left, t_top + t_row_h, t_w, Inches(0.025), DIVIDER)

data_rows = [
    ("In input view",    "~20.5k", "✓  Yes", "True",  "True",  "✗  NO",    ORANGE, BADGE_ORANGE),
    ("Valid — held out", "~20.5k", "✗  No",  "False", "True",  "✓  YES",   GREEN,  BADGE_GREEN),
    ("NaN position",     "~8k",    "✗  No",  "False", "False", "✗  NO",    RED,    BADGE_RED),
]
for i, (pos, cnt, sees, inm, vm, inloss, color, bg) in enumerate(data_rows):
    y = t_top + t_row_h + Inches(0.025) + i * (t_row_h + Inches(0.04))
    add_rect(slide, t_left, y, t_w, t_row_h, bg)
    il_color = GREEN if "YES" in inloss else RED
    for text, w, x, fc in [
        (pos,    col_ws[0], col_xs[0], color),
        (cnt,    col_ws[1], col_xs[1], TEXT_BODY),
        (sees,   col_ws[2], col_xs[2], GREEN if "Yes" in sees else TEXT_LIGHT),
        (inm,    col_ws[3], col_xs[3], ORANGE if inm == "True" else TEXT_LIGHT),
        (vm,     col_ws[4], col_xs[4], GREEN if vm == "True" else RED),
        (inloss, col_ws[5], col_xs[5], il_color),
    ]:
        add_textbox(slide, text, x + Inches(0.08), y + Inches(0.12),
                    w - Inches(0.12), Inches(0.38), font_size=13, color=fc, bold=(fc==il_color and "YES" in text))

# Bottom panels
card(slide, Inches(0.4), Inches(5.72), Inches(6.0), Inches(1.35),
     bg=BADGE_BLUE, border_color=BLUE)
add_textbox(slide, "Loss formula",
            Inches(0.55), Inches(5.79), Inches(5.7), Inches(0.35),
            font_size=13, bold=True, color=BLUE_DARK)
add_textbox(slide,
    "loss = Σ  MSE(pred, target) × recon_mask\n"
    "     ÷  recon_mask.sum()  (true denominator)\n"
    "≈ 42% of 49k positions contribute per step",
    Inches(0.55), Inches(6.14), Inches(5.7), Inches(0.85),
    font_size=13, color=TEXT_BODY)

card(slide, Inches(6.7), Inches(5.72), Inches(6.25), Inches(1.35),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide, "Metrics (PCC, MAE) use same recon_mask",
            Inches(6.85), Inches(5.79), Inches(5.9), Inches(0.35),
            font_size=13, bold=True, color=GREEN_DARK)
add_textbox(slide,
    "pred_values  = pred_betas[recon_mask]\n"
    "true_values  = all_betas[recon_mask]\n"
    "PCC computed only on real, held-out CpGs  →  honest metric",
    Inches(6.85), Inches(6.14), Inches(5.9), Inches(0.85),
    font_size=13, color=TEXT_BODY)


# =============================================================================
# SLIDE 9 — Bug vs Fix
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Bug in Old Code vs. Correct Implementation in LLaMA",
             "Your LLaMA runs are NOT affected — the fix is already in wced_llama.py  ✓")

# Side-by-side comparison (visual, no code)
card(slide, Inches(0.4), Inches(1.38), Inches(6.1), Inches(5.7),
     bg=RGBColor(0xFF,0xF5,0xF5), border_color=RED)
add_rect(slide, Inches(0.4), Inches(1.38), Inches(6.1), Inches(0.06), RED)
add_textbox(slide, "✗  BUG — wced_module.py  (SCBert)",
            Inches(0.55), Inches(1.48), Inches(5.8), Inches(0.42),
            font_size=17, bold=True, color=RED)
badge(slide, "NOT YOUR RUN", Inches(0.55), Inches(1.92), Inches(1.8), Inches(0.28),
      bg=BADGE_RED, fg=RED, font_size=11, bold=True)

bug_points = [
    ("recon_mask = ~input_mask ONLY",
     "valid_mask is in the batch dict but never read"),
    ("Denominator includes NaN positions",
     "~29k NaN positions with target=0 are in the loss denominator"),
    ("75% of gradient is from fake zeros",
     "model is trained on artificial 0.0 values it never actually saw"),
    ("MSE looks artificially low",
     "easy-to-predict zeros pull the loss down — misleading metric"),
    ("PCC computed on all non-input",
     "correlation diluted by 29k positions with target=0"),
]
for i, (title, desc) in enumerate(bug_points):
    y = Inches(2.32 + i * 0.82)
    add_rect(slide, Inches(0.55), y + Inches(0.12), Inches(0.06), Inches(0.5), RED)
    add_textbox(slide, title, Inches(0.72), y + Inches(0.1),
                Inches(5.5), Inches(0.38), font_size=13, bold=True, color=RED)
    add_textbox(slide, desc, Inches(0.72), y + Inches(0.46),
                Inches(5.5), Inches(0.38), font_size=12, color=TEXT_BODY)

card(slide, Inches(7.2), Inches(1.38), Inches(5.75), Inches(5.7),
     bg=RGBColor(0xF0,0xFB,0xF5), border_color=GREEN)
add_rect(slide, Inches(7.2), Inches(1.38), Inches(5.75), Inches(0.06), GREEN)
add_textbox(slide, "✓  CORRECT — wced_llama.py  (LLaMA)",
            Inches(7.35), Inches(1.48), Inches(5.4), Inches(0.42),
            font_size=17, bold=True, color=GREEN)
badge(slide, "YOUR PRETRAIN RUN 44450919", Inches(7.35), Inches(1.92), Inches(2.7), Inches(0.28),
      bg=BADGE_GREEN, fg=GREEN_DARK, font_size=11, bold=True)

fix_points = [
    ("recon_mask = ~input_mask  AND  valid_mask",
     "valid_mask explicitly applied — NaN positions excluded"),
    ("Denominator = recon_mask.sum()",
     "only counts positions that are real AND held-out"),
    ("100% of gradient from real CpGs",
     "every position contributing to loss has a true β measurement"),
    ("MSE is honest",
     "no artificial zeros — metric reflects true reconstruction quality"),
    ("PCC on real held-out only",
     "correlation computed on the same clean mask as the loss"),
]
for i, (title, desc) in enumerate(fix_points):
    y = Inches(2.32 + i * 0.82)
    add_rect(slide, Inches(7.35), y + Inches(0.12), Inches(0.06), Inches(0.5), GREEN)
    add_textbox(slide, title, Inches(7.52), y + Inches(0.1),
                Inches(5.2), Inches(0.38), font_size=13, bold=True, color=GREEN)
    add_textbox(slide, desc, Inches(7.52), y + Inches(0.46),
                Inches(5.2), Inches(0.38), font_size=12, color=TEXT_BODY)

# Impact comparison table
card(slide, Inches(0.4), Inches(7.1), Inches(12.55), Inches(0.22), bg=PANEL_BG)
for i, (label, old_val, new_val) in enumerate([
    ("Real signal:", "25%", "~100%"),
]):
    pass  # handled below

add_textbox(slide,
    "Signal quality:  OLD = 25% real  ·  NEW = ~100% real     "
    "Metric honesty:  OLD = inflated  ·  NEW = correct     "
    "Your run:  ✓ Already correct from day 1",
    Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.2),
    font_size=12, color=GREEN_DARK)


# =============================================================================
# SLIDE 10 — Pretrain Metrics
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Pretrain Results — Run 44450919  (Epoch 99 / 300)",
             "llama-small-all49k-r0.5-w0.0  ·  256D × 4L × 4H  ·  ~5M params  ·  Runtime: ~4.5 days")

# Top metrics
top_m = [
    ("0.971",  "val / pcc\nReconstruction quality",    GREEN,  BADGE_GREEN),
    ("0.043",  "val / mae\nβ error on [0, 1]",          BLUE,   BADGE_BLUE),
    ("0.006",  "val / loss\nRecon MSE",                  BLUE,   BADGE_BLUE),
    ("42.0%",  "valid_pct\nUseful positions / step",    ORANGE, BADGE_ORANGE),
]
for i, (val, label, vc, bg) in enumerate(top_m):
    metric_card(slide, Inches(0.4 + i * 3.15), Inches(1.38),
                Inches(2.9), Inches(1.25), val, label, val_color=vc, bg=bg)

# CLS diagnostics
section_label(slide, "CLS COLLAPSE DIAGNOSTICS  (collapsed = bad)", Inches(0.4), Inches(2.82))
cls_m = [
    ("0.376",  "cls_similarity\n(1.0 = fully collapsed)",   GREEN,  BADGE_GREEN),
    ("0.429",  "cls_variance\n(0.0 = fully collapsed)",      GREEN,  BADGE_GREEN),
    ("0.310",  "pred_std\n(0.0 = predicting mean)",          GREEN,  BADGE_GREEN),
    ("0.570",  "pred_var_ratio\n(1.0 = perfect spread)",     ORANGE, BADGE_ORANGE),
]
for i, (val, label, vc, bg) in enumerate(cls_m):
    metric_card(slide, Inches(0.4 + i * 3.15), Inches(3.12),
                Inches(2.9), Inches(1.25), val, label, val_color=vc, bg=bg)

# Two info panels
card(slide, Inches(0.4), Inches(4.55), Inches(6.1), Inches(2.05),
     bg=PANEL_BG, border_color=BLUE)
section_label(slide, "NaN RATE INFERENCE FROM valid_pct", Inches(0.55), Inches(4.62), BLUE)
add_textbox(slide,
    "Investigation confirms:  NaN rate = 8.69% in pretrain\n"
    "→  ~44.9k of 49k probes measured per sample on average\n"
    "valid_pct ≈ 42% = (0.5 × 44.9k) / 49.15k  ✓ consistent\n"
    "Pretrain corpus is dominated by EPIC-array samples",
    Inches(0.55), Inches(4.98), Inches(5.7), Inches(1.5),
    font_size=13, color=TEXT_BODY)

card(slide, Inches(6.7), Inches(4.55), Inches(6.25), Inches(2.05),
     bg=PANEL_BG, border_color=DIVIDER)
section_label(slide, "TRAINING HEALTH", Inches(6.85), Inches(4.62), BLUE)
health = [
    ("cpg_scale", "3.68  ✓  (init=0.1 → growing healthily)", GREEN),
    ("grad_norm",  "0.0018  (stable convergence)",             BLUE),
    ("train/pcc",  "0.970  (no overfitting)",                  GREEN),
    ("LR @ ep 99", "3.82e-4  (cosine decay on track)",         TEXT_BODY),
]
for i, (k, v, fc) in enumerate(health):
    y = Inches(4.98 + i * 0.38)
    add_textbox(slide, k, Inches(6.85), y, Inches(1.6), Inches(0.35),
                font_size=12, bold=True, color=TEXT_LIGHT)
    add_textbox(slide, v, Inches(8.5), y, Inches(4.35), Inches(0.35),
                font_size=12, color=fc)

# Interpretation bar
card(slide, Inches(0.4), Inches(6.75), Inches(12.55), Inches(0.5),
     bg=BADGE_GREEN, border_color=GREEN)
add_textbox(slide,
    "✓ No CLS collapse   ✓ Reconstruction excellent (PCC 0.971)   "
    "⚠ pred_var_ratio = 0.57 — slight underdispersion (normal for pure reconstruction pretraining)",
    Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.38),
    font_size=13, bold=True, color=GREEN_DARK)


# =============================================================================
# SLIDE 11 — Investigation Results — NaN Statistics (charts)
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Data Investigation Results — NaN Statistics",
             "Job 44584079  ·  Mon Apr 27 10:20–10:28 IDT 2026")

# ── Chart A: NaN % comparison bar ───────────────────────────────────────────
fig, ax = plt.subplots(figsize=(3.8, 2.6))
fig.patch.set_facecolor("white")
datasets = ["Pretrain\n(169k samples)", "Fine-tune\n(11.5k samples)"]
nan_pcts = [8.69, 60.11]
colors_bar = ["#007AC2", "#C03030"]
bars = ax.bar(datasets, nan_pcts, color=colors_bar, width=0.5, edgecolor="white", linewidth=1.5)
for bar, pct in zip(bars, nan_pcts):
    ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.5,
            f"{pct}%", ha="center", va="bottom", fontsize=14, fontweight="bold",
            color=bar.get_facecolor())
ax.set_ylabel("NaN %", fontsize=11, color="#3A4A5C")
ax.set_ylim(0, 75)
ax.set_title("Overall NaN Rate", fontsize=13, fontweight="bold", color="#1A263A", pad=8)
ax.tick_params(colors="#3A4A5C", labelsize=10)
ax.spines[["top", "right"]].set_visible(False)
ax.spines[["left", "bottom"]].set_color("#CCD9E8")
ax.yaxis.grid(True, color="#EEF2F8", linewidth=0.8)
ax.set_axisbelow(True)
fig.tight_layout()
fig_to_slide(slide, fig, Inches(0.4), Inches(1.42), Inches(4.2), Inches(2.9))

# ── Chart B: Pretrain — % samples by NaN threshold ──────────────────────────
fig, ax = plt.subplots(figsize=(4.0, 2.6))
fig.patch.set_facecolor("white")
thresholds = ["≥1%", "≥5%", "≥10%", "≥25%", "≥50%"]
pcts = [54.2, 27.5, 20.9, 11.2, 10.0]
bar_colors = ["#007AC2"] * 5
bars2 = ax.barh(thresholds[::-1], pcts[::-1], color=bar_colors, edgecolor="white")
for bar, pct in zip(bars2, pcts[::-1]):
    ax.text(pct + 0.5, bar.get_y() + bar.get_height() / 2,
            f"{pct}%", va="center", fontsize=10, color="#007AC2", fontweight="bold")
ax.set_xlabel("% of samples", fontsize=10, color="#3A4A5C")
ax.set_xlim(0, 70)
ax.set_title("Pretrain: % samples with NaN ≥ threshold", fontsize=11,
             fontweight="bold", color="#1A263A", pad=6)
ax.tick_params(colors="#3A4A5C", labelsize=10)
ax.spines[["top", "right"]].set_visible(False)
ax.spines[["left", "bottom"]].set_color("#CCD9E8")
ax.xaxis.grid(True, color="#EEF2F8", linewidth=0.8)
ax.set_axisbelow(True)
fig.tight_layout()
fig_to_slide(slide, fig, Inches(4.75), Inches(1.42), Inches(4.3), Inches(2.9))

# ── Chart C: Fine-tune CpG coverage breakdown ───────────────────────────────
fig, ax = plt.subplots(figsize=(3.6, 2.6))
fig.patch.set_facecolor("white")
labels = ["Valid CpGs\n(0 NaN)\n19,608", "100% NaN CpGs\n(always missing)\n29,548"]
sizes = [19608, 29548]
pie_colors = ["#008F6A", "#C03030"]
wedges, texts, autotexts = ax.pie(
    sizes, labels=labels, colors=pie_colors,
    autopct="%1.1f%%", startangle=90,
    textprops={"fontsize": 9, "color": "#1A263A"},
    wedgeprops={"edgecolor": "white", "linewidth": 2}
)
for at in autotexts:
    at.set_fontsize(11)
    at.set_fontweight("bold")
    at.set_color("white")
ax.set_title("Fine-tune: CpG Coverage\n(out of 49,156 total)", fontsize=11,
             fontweight="bold", color="#1A263A", pad=4)
fig.tight_layout()
fig_to_slide(slide, fig, Inches(9.2), Inches(1.42), Inches(3.9), Inches(2.9))

# ── Stats boxes below charts ─────────────────────────────────────────────────
for i, (label, val, vc, bg_c) in enumerate([
    ("Pretrain NaN %",          "8.69%",        BLUE,   BADGE_BLUE),
    ("Fine-tune NaN %",         "60.11%",       RED,    BADGE_RED),
    ("Fine-tune fully-missing\nCpGs (100% NaN)","29,548 / 49,156", RED,   BADGE_RED),
    ("Fine-tune splits\ntrain / val / test",    "5,461 / 1,366 / 4,626", BLUE, BADGE_BLUE),
]):
    lx = Inches(0.4 + i * 3.22)
    card(slide, lx, Inches(4.48), Inches(3.05), Inches(1.18), bg=bg_c, border_color=None)
    add_textbox(slide, val, lx, Inches(4.54), Inches(3.05), Inches(0.55),
                font_size=22, bold=True, color=vc, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, lx, Inches(5.08), Inches(3.05), Inches(0.52),
                font_size=11, color=TEXT_BODY, align=PP_ALIGN.CENTER)

# Key insight bar
card(slide, Inches(0.4), Inches(5.82), Inches(12.55), Inches(0.78),
     bg=BADGE_BLUE, border_color=BLUE)
add_textbox(slide,
    "Key insight:  Fine-tune = pure 450k-array dataset — same 29,548 CpGs are 100% NaN "
    "across every single sample.  "
    "Pretrain = mostly EPIC-array (NaN only 8.69%).  "
    "The WCEDCollator valid_mask handles both correctly.",
    Inches(0.55), Inches(5.88), Inches(12.2), Inches(0.65),
    font_size=13, bold=True, color=BLUE_DARK)


# =============================================================================
# SLIDE 12 — Investigation Results — Beta Distribution (charts)
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Data Investigation Results — β-Value Distribution",
             "Valid values only  ·  NaN excluded  ·  Pretrain: 7.6B values  ·  Fine-tune: 225M values")

# ── Chart A: Overlaid CDF / percentile curve ─────────────────────────────────
percs_x  = [0, 1, 5, 25, 50, 75, 95, 99, 100]
pre_vals = [-0.3389, 0.0100, 0.0236, 0.0700, 0.2860, 0.6679, 0.9134, 0.9620, 1.1030]
ft_vals  = [ 0.0000, 0.0086, 0.0188, 0.0400, 0.0825, 0.4881, 0.9050, 0.9498, 1.0000]

fig, ax = plt.subplots(figsize=(6.2, 4.0))
fig.patch.set_facecolor("white")
ax.plot(pre_vals, percs_x, color="#007AC2", linewidth=2.5, marker="o",
        markersize=6, label="Pretrain (169k samples)")
ax.plot(ft_vals,  percs_x, color="#008F6A", linewidth=2.5, marker="s",
        markersize=6, label="Fine-tune (11.5k samples)", linestyle="--")
ax.axvline(0, color="#888", linewidth=0.8, linestyle=":")
ax.axvline(1, color="#888", linewidth=0.8, linestyle=":")
ax.fill_betweenx(percs_x, pre_vals, ft_vals, alpha=0.08, color="#007AC2")
ax.set_xlabel("β-value", fontsize=12, color="#3A4A5C")
ax.set_ylabel("Percentile", fontsize=12, color="#3A4A5C")
ax.set_title("β-Value CDF: Pretrain vs Fine-tune", fontsize=13,
             fontweight="bold", color="#1A263A", pad=8)
ax.legend(fontsize=10, framealpha=0.9, loc="upper left")
ax.tick_params(colors="#3A4A5C", labelsize=10)
ax.spines[["top", "right"]].set_visible(False)
ax.spines[["left", "bottom"]].set_color("#CCD9E8")
ax.grid(True, color="#EEF2F8", linewidth=0.8)
ax.set_axisbelow(True)
# median lines
ax.axhline(50, color="#AAA", linewidth=0.8, linestyle=":")
ax.annotate("Pretrain median\n0.286", xy=(0.286, 50), xytext=(0.35, 60),
            fontsize=9, color="#007AC2",
            arrowprops=dict(arrowstyle="->", color="#007AC2", lw=1))
ax.annotate("Fine-tune median\n0.083", xy=(0.083, 50), xytext=(0.2, 35),
            fontsize=9, color="#008F6A",
            arrowprops=dict(arrowstyle="->", color="#008F6A", lw=1))
fig.tight_layout()
fig_to_slide(slide, fig, Inches(0.4), Inches(1.4), Inches(6.8), Inches(4.4))

# ── Chart B: bar chart of key percentiles side by side ───────────────────────
fig, ax = plt.subplots(figsize=(5.4, 4.0))
fig.patch.set_facecolor("white")
perc_labels = ["1%", "5%", "25%", "50%", "75%", "95%", "99%"]
pre_v = [0.0100, 0.0236, 0.0700, 0.2860, 0.6679, 0.9134, 0.9620]
ft_v  = [0.0086, 0.0188, 0.0400, 0.0825, 0.4881, 0.9050, 0.9498]
x = np.arange(len(perc_labels))
w = 0.38
b1 = ax.bar(x - w/2, pre_v, w, label="Pretrain", color="#007AC2",
            edgecolor="white", linewidth=1)
b2 = ax.bar(x + w/2, ft_v,  w, label="Fine-tune", color="#008F6A",
            edgecolor="white", linewidth=1)
ax.set_xticks(x)
ax.set_xticklabels(perc_labels, fontsize=10, color="#3A4A5C")
ax.set_ylabel("β-value", fontsize=11, color="#3A4A5C")
ax.set_title("β-Value Percentiles Compared", fontsize=12,
             fontweight="bold", color="#1A263A", pad=8)
ax.legend(fontsize=10, framealpha=0.9)
ax.tick_params(colors="#3A4A5C", labelsize=10)
ax.spines[["top", "right"]].set_visible(False)
ax.spines[["left", "bottom"]].set_color("#CCD9E8")
ax.yaxis.grid(True, color="#EEF2F8", linewidth=0.8)
ax.set_axisbelow(True)
fig.tight_layout()
fig_to_slide(slide, fig, Inches(7.4), Inches(1.4), Inches(5.6), Inches(4.4))

# ── Key stat boxes ────────────────────────────────────────────────────────────
for i, (label, val, vc, bg_c) in enumerate([
    ("Pretrain median β",      "0.286",   BLUE,   BADGE_BLUE),
    ("Fine-tune median β",     "0.083",   GREEN,  BADGE_GREEN),
    ("Pretrain out-of-range",  "248,776 ⚠", ORANGE, BADGE_ORANGE),
    ("Fine-tune out-of-range", "0  ✓",    GREEN,  BADGE_GREEN),
]):
    lx = Inches(0.4 + i * 3.22)
    card(slide, lx, Inches(5.98), Inches(3.05), Inches(0.92), bg=bg_c, border_color=None)
    add_textbox(slide, val, lx, Inches(6.0), Inches(3.05), Inches(0.48),
                font_size=20, bold=True, color=vc, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, lx, Inches(6.48), Inches(3.05), Inches(0.36),
                font_size=11, color=TEXT_BODY, align=PP_ALIGN.CENTER)

card(slide, Inches(0.4), Inches(7.0), Inches(12.55), Inches(0.38),
     bg=BADGE_ORANGE, border_color=ORANGE)
add_textbox(slide,
    "Distribution shift:  Pretrain is bimodal (unmethylated 0–0.2  +  methylated 0.6–1.0)  ·  "
    "Fine-tune is more heavily unmethylated (median 0.083)  ·  "
    "Pretrain out-of-range due to raw Illumina signal (harmless — model learns from valid range)",
    Inches(0.55), Inches(7.04), Inches(12.2), Inches(0.3),
    font_size=12, color=ORANGE_DARK)


# =============================================================================
# SLIDE 13 — Summary
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
slide_header(slide, "Summary",
             "The full NaN handling chain is correct in both pretrain and fine-tune LLaMA runs")

card(slide, Inches(0.4), Inches(1.38), Inches(5.9), Inches(5.75),
     bg=PANEL_BG, border_color=BLUE)
section_label(slide, "PIPELINE — WHAT EACH FILE DOES", Inches(0.55), Inches(1.48), BLUE)
items = [
    ("Step 1  ·  MethylationDataset",    "Reads β  ·  computes valid_mask = isfinite(β)  ·  NaN preserved", BLUE),
    ("Step 2a  ·  WCEDCollator",         "NaN → 0 in all_betas  ·  valid_mask stored in batch dict",       ORANGE),
    ("Step 2b  ·  WCEDCollator",         "Input views: only sample from valid_indices  ·  NaN never tokenised", PURPLE),
    ("Step 3  ·  WCEDLlamaModule",       "recon_mask = ~input_mask & valid_mask  ·  honest denominator",    GREEN),
]
for i, (title, desc, color) in enumerate(items):
    y = Inches(1.9 + i * 1.18)
    add_rect(slide, Inches(0.55), y, Inches(0.06), Inches(1.0), color)
    add_textbox(slide, title, Inches(0.75), y + Inches(0.06),
                Inches(5.2), Inches(0.38), font_size=13, bold=True, color=color)
    add_textbox(slide, desc, Inches(0.75), y + Inches(0.48),
                Inches(5.2), Inches(0.45), font_size=12, color=TEXT_BODY)

card(slide, Inches(6.75), Inches(1.38), Inches(6.2), Inches(2.7),
     bg=BADGE_GREEN, border_color=GREEN)
section_label(slide, "YOUR LLAMA RUNS — STATUS", Inches(6.9), Inches(1.48), GREEN_DARK)
for i, text in enumerate([
    "✓  Pretrain 44450919: correct NaN handling",
    "✓  Fine-tune 44574410: recon_weight=0.0 → NaN irrelevant",
    "✓  wced_llama.py had the fix from the very start",
    "✓  Bug exists only in old SCBert code (not your runs)",
]):
    add_textbox(slide, text, Inches(6.9), Inches(1.9 + i * 0.48),
                Inches(5.8), Inches(0.42), font_size=13, color=GREEN_DARK)

card(slide, Inches(6.75), Inches(4.22), Inches(6.2), Inches(2.9),
     bg=BADGE_ORANGE, border_color=ORANGE)
section_label(slide, "PRETRAIN QUALITY  (EPOCH 99)", Inches(6.9), Inches(4.32), ORANGE_DARK)
checks = [
    ("PCC = 0.971",          "excellent reconstruction",      GREEN),
    ("MAE = 0.043",          "4.3% error on [0,1]",           GREEN),
    ("cls_similarity = 0.376", "no CLS collapse",             GREEN),
    ("cls_variance = 0.429",  "healthy representation diversity", GREEN),
    ("pred_var_ratio = 0.570", "slight underdispersion — expected", ORANGE),
    ("valid_pct ≈ 42%",      "NaN rate ≈ 16% in pretrain corpus", BLUE),
]
for i, (k, v, fc) in enumerate(checks):
    y = Inches(4.72 + i * 0.38)
    add_textbox(slide, k, Inches(6.9), y, Inches(2.4), Inches(0.35),
                font_size=12, bold=True, color=fc)
    add_textbox(slide, v, Inches(9.35), y, Inches(3.45), Inches(0.35),
                font_size=12, color=TEXT_BODY)

card(slide, Inches(0.4), Inches(7.2), Inches(12.55), Inches(0.22), bg=PANEL_BG)
add_textbox(slide,
    "Next: paste investigate_data.sh results into slides 11–12  ·  "
    "Monitor finetune val/r2 to assess if reconstruction pretraining transferred to age prediction",
    Inches(0.55), Inches(7.22), Inches(12.2), Inches(0.2),
    font_size=12, color=TEXT_LIGHT, italic=True)


# =============================================================================
# Save
# =============================================================================
out_path = "NaN_Handling_MethylLlama.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
